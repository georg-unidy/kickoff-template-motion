import json, os, io, requests
from pptx import Presentation
from pptx.util import Emu
from pptx.dml.color import RGBColor

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
TEMPLATE_DE = os.path.join(BASE_DIR, "Mein_Holstein_Kiel_SSO_Kick-off.pptx")
TEMPLATE_EN = os.path.join(BASE_DIR, "SFL_EN_Template.pptx")

# Brand accent colors to replace
DE_COLOR      = "005397"   # HK blue — text + setup boxes
EN_COLOR      = "003DA5"   # SFL blue — text
EN_FILL_COLOR  = "007699"   # SFL teal — setup boxes (different from text color!)
EN_COVER_TEXT  = "0055FF"   # SFL cover headline color (Brand name)

# Logo: top-right, 6.61 x 6.61 cm (same as HK)
LOGO_LEFT = Emu(int(24.82 * 914400 / 2.54))
LOGO_TOP  = Emu(int(0.85  * 914400 / 2.54))
LOGO_SIZE = Emu(int(6.61  * 914400 / 2.54))

AGENDA = [
    "Intro Unidy",
    "Intro Instance & Admin",
    "Setup & Versions",
    "Intro Customer Page",
]

# Slide indices (0-based) — identical for both templates after trimming
IDX_COVER   = 0
IDX_SETUP   = 7   # slide 8


def replace_run_text(shape, old, new):
    if not shape.has_text_frame:
        return False
    changed = False
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if old in run.text:
                run.text = run.text.replace(old, new)
                changed = True
    return changed

def hex_to_rgb(h):
    h = h.lstrip("#").upper()
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

def replace_fill_color(prs, old, new):
    rgb = hex_to_rgb(new)
    old = old.upper().lstrip("#")
    for slide in prs.slides:
        for shape in slide.shapes:
            try:
                if str(shape.fill.fore_color.rgb).upper() == old:
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = rgb
            except:
                pass

def replace_text_color(prs, old, new):
    rgb = hex_to_rgb(new)
    old = old.upper().lstrip("#")
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    try:
                        if str(run.font.color.rgb).upper() == old:
                            run.font.color.rgb = rgb
                    except:
                        pass

def is_agenda_shape(shape):
    if not shape.has_text_frame:
        return False
    paras = [p for p in shape.text_frame.paragraphs
             if ''.join(r.text for r in p.runs).strip()]
    if len(paras) != 4:
        return False
    combined = ' '.join(''.join(r.text for r in p.runs) for p in paras)
    return 'Intro' in combined and ('Setup' in combined or 'Versions' in combined)

def set_agenda_shape(shape, bold_idx):
    paras = [p for p in shape.text_frame.paragraphs
             if ''.join(r.text for r in p.runs).strip()]
    for i, para in enumerate(paras):
        new_text = AGENDA[i] if i < len(AGENDA) else ""
        for run in para.runs:
            run.text = ""
        if para.runs:
            para.runs[0].text = new_text
            para.runs[0].font.bold = (i == bold_idx)

def replace_cover_text_colors(slide, new_hex):
    """
    Slide 1 EN: replace both explicit 0055FF runs AND accent2 scheme runs
    with the CI color so headline + date + subtitle all match.
    """
    from pptx.oxml.ns import qn
    new_rgb = hex_to_rgb(new_hex)
    ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                rpr = run._r.find(qn("a:rPr"))
                if rpr is None:
                    continue
                # Case 1: explicit srgbClr (e.g. 0055FF)
                srgb = rpr.find(f".//{{{ns}}}srgbClr")
                if srgb is not None:
                    srgb.set("val", new_hex.upper().lstrip("#"))
                # Case 2: schemeClr accent2 — replace with explicit srgb
                scheme = rpr.find(f".//{{{ns}}}schemeClr")
                if scheme is not None and scheme.get("val") == "accent2":
                    solid = rpr.find(f".//{{{ns}}}solidFill")
                    if solid is not None:
                        solid.remove(scheme)
                        from lxml import etree
                        srgb_new = etree.SubElement(solid, f"{{{ns}}}srgbClr")
                        srgb_new.set("val", new_hex.upper().lstrip("#"))

def replace_center_crest(slide, logo_url, is_en):
    """Replace the center crest on the setup diagram with the customer logo."""
    # HK: Shape 20 is center crest; SFL EN: Shape 15
    # Identify by position: left ~11-12cm, top ~9-11cm, width < 3cm
    target = None
    for shape in slide.shapes:
        if shape.shape_type != 13:
            continue
        left = shape.left / 914400 * 2.54
        top  = shape.top  / 914400 * 2.54
        w    = shape.width / 914400 * 2.54
        if 9 < left < 13 and 8 < top < 12 and w < 3:
            target = shape
            break
    if target is None:
        return

    # Save position and size
    left   = target.left
    top    = target.top
    width  = target.width
    height = target.height

    # Remove old crest
    el = target._element
    el.getparent().remove(el)

    # Insert new logo at same position
    import requests as req, io
    r = req.get(logo_url, timeout=15)
    r.raise_for_status()
    slide.shapes.add_picture(io.BytesIO(r.content), left, top, width, height)

def remove_picture(slide, index):
    pics = [s for s in slide.shapes if s.shape_type == 13]
    if index < len(pics):
        el = pics[index]._element
        el.getparent().remove(el)

def insert_logo(slide, logo_url):
    r = requests.get(logo_url, timeout=15)
    r.raise_for_status()
    slide.shapes.add_picture(io.BytesIO(r.content), LOGO_LEFT, LOGO_TOP, LOGO_SIZE, LOGO_SIZE)


def generate(config: dict, output_path: str):
    brand        = config.get("brand_name", "Brand")
    date         = config.get("kickoff_date", "TBD")
    language     = config.get("language", "de").strip().lower()
    csm_email    = config.get("csm_email", "service@unidy.de")
    csm_name     = config.get("csm_name", "")
    csm_phone    = config.get("csm_phone", "")
    goals        = config.get("goals", ["", "", "", ""])
    integrations = config.get("integrations_v1", [])
    go_live_date = config.get("go_live_date", "")
    notion_url   = config.get("notion_url", "")
    logo_url     = config.get("logo_url", "")
    ci_primary   = config.get("ci_primary_color", "").lstrip("#").upper()
    ci_text      = config.get("ci_text_color", "").lstrip("#").upper()

    while len(goals) < 4:
        goals.append("")

    is_en = language in ("en", "english", "englisch")
    template_path = TEMPLATE_EN if is_en else TEMPLATE_DE
    base_color    = EN_COLOR if is_en else DE_COLOR

    prs = Presentation(template_path)

    parts = date.split("/")
    dd   = parts[0] if len(parts) > 0 else ""
    mm   = parts[1] if len(parts) > 1 else ""
    yyyy = parts[2] if len(parts) > 2 else ""

    # ── Cover (Slide 1) ───────────────────────────────────────────────────────
    slide1 = prs.slides[IDX_COVER]

    if is_en:
        for shape in slide1.shapes:
            replace_run_text(shape, "SFL ID", f"{brand} ID")
            replace_run_text(shape, "SFL", brand)
            replace_run_text(shape, "15/12/2025", f"{dd}/{mm}/{yyyy}")
    else:
        for shape in slide1.shapes:
            # HK has "Mein Holstein Kiel" without " ID"
            replace_run_text(shape, "Mein Holstein Kiel", f"{brand} ID")
            # Date split: '03' '/06' '/2026'
            replace_run_text(shape, "03", dd)
            replace_run_text(shape, "/06", f"/{mm}")
            replace_run_text(shape, "/2026", f"/{yyyy}")

    if logo_url:
        remove_picture(slide1, 1)   # index 1 = customer logo (index 0 = Unidy logo)
        insert_logo(slide1, logo_url)

    # ── Agenda (all agenda slides) ────────────────────────────────────────────
    agenda_count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if is_agenda_shape(shape):
                set_agenda_shape(shape, bold_idx=agenda_count % 4)
                agenda_count += 1
                break

    # ── Setup slide (Slide 8) ─────────────────────────────────────────────────
    if logo_url:
        replace_center_crest(prs.slides[IDX_SETUP], logo_url, is_en)
    for shape in prs.slides[IDX_SETUP].shapes:
        if is_en:
            replace_run_text(shape, "SFL Setup", f"{brand} Setup")
            replace_run_text(shape, "SFL", brand)
        else:
            replace_run_text(shape, "Mein Holstein Kiel SSO Setup", f"{brand} SSO Setup")
            replace_run_text(shape, "Mein Holstein Kiel", brand)

    # ── Goals (anywhere in deck) ──────────────────────────────────────────────
    for slide in prs.slides:
        for shape in slide.shapes:
            for i, goal in enumerate(goals[:4], 1):
                if goal:
                    replace_run_text(shape, f"Goal {i}", goal)
            replace_run_text(shape, "Hier noch JGS Goals einbauen", "")
            replace_run_text(shape,
                "https://www.notion.so/unidy-gmbh/Customer-Outcomes-and-Customer-Personas-31e5a77a0d338015b860eca1bda9a136", "")

    # ── Next Steps ────────────────────────────────────────────────────────────
    for slide in prs.slides:
        for shape in slide.shapes:
            if integrations:
                replace_run_text(shape,
                    "Kunden kontaktiert alle Integrationspartner mit ",
                    f"Integrationen V1: {' · '.join(integrations[:3])} - Partner kontaktieren mit ")
            if go_live_date:
                replace_run_text(shape,
                    "Go-Live Checklist on Customer Page with Customer To Dos",
                    f"Go-Live Checklist - Ziel: {go_live_date}")
            if notion_url:
                replace_run_text(shape, "Erwähnung Notion Page", f"Notion: {notion_url}")

    # ── CI Colors ─────────────────────────────────────────────────────────────
    if ci_primary:
        replace_fill_color(prs, base_color, ci_primary)
        replace_text_color(prs, base_color, ci_primary)
        if is_en:
            # EN setup boxes use a different base color
            replace_fill_color(prs, EN_FILL_COLOR, ci_primary)
            # EN cover uses 0055FF + accent2 scheme — override both
            replace_text_color(prs, EN_COVER_TEXT, ci_primary)
            replace_cover_text_colors(prs.slides[IDX_COVER], ci_primary)
    if ci_text and ci_text != ci_primary:
        replace_text_color(prs, ci_primary if ci_primary else base_color, ci_text)

    prs.save(output_path)


if __name__ == "__main__":
    import argparse
    parser = argparse.ArgumentParser()
    parser.add_argument("--config", required=True)
    parser.add_argument("--output", required=True)
    args = parser.parse_args()
    with open(args.config) as f:
        config = json.load(f)
    generate(config, args.output)
